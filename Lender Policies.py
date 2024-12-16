from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Key Lender Policies and Regulations"
subtitle.text = "A Summary for Credit Analysts"

# Slide 2: Summary of Key Lender Policies
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
title.text = "Summary of Key Lender Policies"

# Add bullet points summarizing lender policies
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Loan-to-Value (LTV) Ratios - Typically capped at 80% to reduce risk\n" \
               "• Debt-to-Income (DTI) Ratios - Most lenders cap DTI at 43% to ensure repayment ability\n" \
               "• Credit Scoring - FICO scores generally used to assess risk\n" \
               "• Interest Rate Structures - Fixed vs. Adjustable rates\n" \
               "• Government Programs - FHA, VA loans help first-time buyers"

# Slide 3: Insights into Lender Policies
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_3.shapes.title
title.text = "Insights into Lender Policies"

# Add insights
content = slide_3.shapes.placeholders[1].text_frame
content.text = "• Lenders typically prefer LTV below 80%, reducing their risk of default.\n" \
               "• Lower DTI ratios indicate better loan repayment ability, increasing approval chances.\n" \
               "• Higher credit scores lead to better loan terms and lower interest rates.\n" \
               "• First-time homebuyer programs can offer lower down payments and better rates."

# Slide 4: Visual 1 - Loan-to-Value (LTV) Ratio Chart
slide_4 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for the chart
title = slide_4.shapes.title
title.text = "Loan-to-Value (LTV) Ratios"

# Generate the LTV chart (example data)
ltv_data = [80, 85, 90, 95]
ltv_labels = ['80% LTV', '85% LTV', '90% LTV', '95% LTV']
plt.bar(ltv_labels, ltv_data, color='skyblue')
plt.title('Loan-to-Value (LTV) Ratios for Different Loan Types')
plt.xlabel('LTV Ratio')
plt.ylabel('Percentage')
plt.tight_layout()

# Save the chart as an image (use relative path)
chart_path = 'ltv_ratio_chart.png'
plt.savefig(chart_path)
plt.close()

# Add the chart to the slide
slide_4.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8))

# Slide 5: Visual 2 - Debt-to-Income (DTI) Ratio Chart
slide_5 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for the chart
title = slide_5.shapes.title
title.text = "Debt-to-Income (DTI) Ratios"

# Generate the DTI chart (example data)
dti_data = [35, 40, 43, 50]
dti_labels = ['35% DTI', '40% DTI', '43% DTI', '50% DTI']
plt.bar(dti_labels, dti_data, color='lightcoral')
plt.title('Debt-to-Income (DTI) Ratios for Different Loan Applications')
plt.xlabel('DTI Ratio')
plt.ylabel('Percentage')
plt.tight_layout()

# Save the chart as an image (use relative path)
dti_chart_path = 'dti_ratio_chart.png'
plt.savefig(dti_chart_path)
plt.close()

# Add the chart to the slide
slide_5.shapes.add_picture(dti_chart_path, Inches(1), Inches(1), width=Inches(8))

# Save the PowerPoint presentation
file_path = 'Lender_Policies_Summary_Presentation.pptx'
prs.save(file_path)

print(f"PowerPoint file saved as {file_path}")
