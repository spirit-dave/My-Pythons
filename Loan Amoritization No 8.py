from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

# Example: Loan details
loan_amount = 300000  # Loan amount in dollars
annual_interest_rate = 0.05  # Annual interest rate (5%)
loan_term_years = 30  # Loan term in years

# Monthly interest rate
monthly_interest_rate = annual_interest_rate / 12

# Number of monthly payments
num_payments = loan_term_years * 12

# Calculate monthly payment using the formula for fixed-rate mortgages
monthly_payment = loan_amount * (monthly_interest_rate * (1 + monthly_interest_rate) ** num_payments) / ((1 + monthly_interest_rate) ** num_payments - 1)

# Generate the amortization schedule
def generate_amortization_schedule(loan_amount, monthly_payment, monthly_interest_rate, num_payments):
    schedule = []
    remaining_balance = loan_amount
    for month in range(1, num_payments + 1):
        interest_payment = remaining_balance * monthly_interest_rate
        principal_payment = monthly_payment - interest_payment
        remaining_balance -= principal_payment
        schedule.append([month, principal_payment, interest_payment, remaining_balance])
    return schedule

# Generate amortization schedule
amortization_schedule = generate_amortization_schedule(loan_amount, monthly_payment, monthly_interest_rate, num_payments)

# Extract principal payments and interest payments for visualization
months = [entry[0] for entry in amortization_schedule]
principal_payments = [entry[1] for entry in amortization_schedule]
interest_payments = [entry[2] for entry in amortization_schedule]

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Loan Amortization Report"
subtitle.text = "Detailed Breakdown of Loan Repayments"

# Slide 2: Financial Data and Insights
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
title.text = "Loan Amortization Details"

# Add bullet points summarizing the loan details and key insights
content = slide_2.shapes.placeholders[1].text_frame
content.text = f"• Loan Amount: ${loan_amount:,}\n" \
               f"• Annual Interest Rate: {annual_interest_rate*100}%\n" \
               f"• Loan Term: {loan_term_years} years\n" \
               f"• Monthly Payment: ${monthly_payment:,.2f}"

# Slide 3: Visualizing Amortization Schedule (Principal vs Interest Payments)
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for chart
title = slide_3.shapes.title
title.text = "Amortization Schedule (Principal vs Interest Payments)"

# Plotting Principal vs Interest Payments
plt.figure(figsize=(8, 6))
plt.plot(months, principal_payments, label='Principal Payment', color='green', linewidth=2)
plt.plot(months, interest_payments, label='Interest Payment', color='red', linewidth=2)
plt.title('Amortization Schedule: Principal vs Interest Payments')
plt.xlabel('Month')
plt.ylabel('Payment Amount ($)')
plt.legend(loc='upper right')
plt.grid(True)

# Save the chart as an image
chart_path = 'amortization_schedule_chart.png'
plt.savefig(chart_path)
plt.close()

# Add the chart to the slide
slide_3.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8))

# Save the PowerPoint presentation
file_path = 'Loan_Amortization_Presentation.pptx'
prs.save(file_path)

print(f"PowerPoint presentation saved as {file_path}")
