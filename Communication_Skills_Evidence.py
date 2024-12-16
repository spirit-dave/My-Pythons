from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Communication Skills Evidence"
subtitle.text = "Summarizing Complex Financial Data Clearly and Effectively"

# Slide 2: Financial Data Summary
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
title.text = "Financial Data Summary"

# Add financial data points to the slide
content = slide_2.shapes.placeholders[1].text_frame
content.text = "Key Metrics:\n" \
               "- Total Borrowing Capacity: $500,000\n" \
               "- Monthly Repayment: $2,200\n" \
               "- Debt-to-Income Ratio: 30%\n" \
               "- Interest Rate: 4.5%\n" \
               "- Loan Term: 30 years"

# Slide 3: Explanation of Data
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_3.shapes.title
title.text = "Explanation of Data"

# Add explanation text
content = slide_3.shapes.placeholders[1].text_frame
content.text = "The borrower’s monthly repayment is calculated based on the loan amount, term, and interest rate. " \
               "The debt-to-income ratio shows the proportion of monthly income used to repay debts, impacting loan eligibility."

# Slide 4: Visual Summary (Placeholder for graphs)
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_4.shapes.title
title.text = "Visual Summary (Graphical Representation)"

# Add placeholder for the visual summary
content = slide_4.shapes.placeholders[1].text_frame
content.text = "In a presentation, the financial data would be visualized using pie charts and bar graphs to highlight key trends:\n" \
               "- Monthly repayment vs. income\n" \
               "- Distribution of assets vs. liabilities"

# Slide 5: Conclusion
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_5.shapes.title
title.text = "Conclusion"

# Add conclusion text
content = slide_5.shapes.placeholders[1].text_frame
content.text = "This presentation demonstrates how financial data can be effectively communicated using both text and visuals. " \
               "Clear communication ensures informed decision-making from stakeholders."

# Save the presentation
file_path = "Communication_Skills_Evidence.pptx"
prs.save(file_path)

print(f"Presentation saved as {file_path}")
