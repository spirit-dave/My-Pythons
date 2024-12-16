from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Example Data (e.g., Loan Portfolio Breakdown)
loan_portfolio = {
    'Home Loans': 500000,
    'Car Loans': 200000,
    'Personal Loans': 100000,
    'Student Loans': 150000,
    'Credit Cards': 50000
}

# Calculate total portfolio and percentage breakdown
total_portfolio = sum(loan_portfolio.values())
portfolio_percentages = {key: (value / total_portfolio) * 100 for key, value in loan_portfolio.items()}

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Loan Portfolio Breakdown"
subtitle.text = "Summary of Portfolio Composition and Insights"

# Slide 2: Financial Data and Insights
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
title.text = "Loan Portfolio Breakdown"

# Add bullet points summarizing the loan portfolio
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• The portfolio consists of 5 main loan types.\n" \
               "• Home loans make up the largest portion (50%) of the portfolio.\n" \
               "• Credit card debt represents a smaller portion (10%) of the overall portfolio.\n" \
               "• The portfolio is well-diversified across different loan types."

# Slide 3: Visualizing Loan Portfolio Breakdown (Pie Chart)
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for chart
title = slide_3.shapes.title
title.text = "Portfolio Breakdown by Loan Type"

# Plotting Pie Chart for Loan Portfolio
plt.figure(figsize=(8, 6))
plt.pie(portfolio_percentages.values(), labels=portfolio_percentages.keys(), autopct='%1.1f%%', startangle=140, colors=['#ff9999','#66b3ff','#99ff99','#ffcc99','#c2c2f0'])
plt.title('Loan Portfolio Breakdown')
plt.axis('equal')  # Equal aspect ratio ensures that pie chart is circular.

# Save the pie chart as an image
chart_path = 'loan_portfolio_pie_chart.png'
plt.savefig(chart_path)
plt.close()

# Add the pie chart to the slide
slide_3.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8))

# Save the PowerPoint presentation
file_path = 'Loan_Portfolio_Presentation.pptx'
prs.save(file_path)

print(f"PowerPoint presentation saved as {file_path}")
