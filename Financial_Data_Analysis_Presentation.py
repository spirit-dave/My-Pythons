from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt

# Data (Example: Analyzing the Borrowing Capacity based on various factors)
loan_amounts = [300000, 350000, 400000, 450000, 500000]  # Loan amounts in dollars
interest_rates = [3.5, 3.7, 4.0, 4.2, 4.5]  # Interest rates in percentage
loan_terms = [15, 20, 25, 30, 35]  # Loan terms in years

# Calculate monthly payments using the formula for loan amortization
def calculate_monthly_payment(principal, annual_rate, term_years):
    monthly_rate = annual_rate / 100 / 12
    num_payments = term_years * 12
    payment = principal * (monthly_rate * (1 + monthly_rate) ** num_payments) / ((1 + monthly_rate) ** num_payments - 1)
    return payment

# Generate monthly payments for each loan configuration
monthly_payments = [
    [calculate_monthly_payment(loan, rate, term) for loan, rate, term in zip(loan_amounts, interest_rates, loan_terms)]
]

# Create a PowerPoint presentation
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Financial Data Analysis: Loan Borrowing Capacity"
subtitle.text = "A Summary Report of Loan Payments"

# Slide 2: Summary of Financial Data and Insights
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
title.text = "Financial Insights on Borrowing Capacity"

# Add bullet points summarizing the financial analysis
content = slide_2.shapes.placeholders[1].text_frame
content.text = "• Loan Amounts range from $300,000 to $500,000.\n" \
               "• Interest rates vary from 3.5% to 4.5%, depending on the loan terms.\n" \
               "• Monthly payments increase with higher loan amounts and longer loan terms.\n" \
               "• Borrowers should consider their ability to handle larger monthly payments for higher loan amounts and longer terms."

# Slide 3: Visualizing Borrowing Capacity Data
slide_3 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for chart
title = slide_3.shapes.title
title.text = "Monthly Payments for Various Loan Terms"

# Plotting the data (monthly payments)
plt.figure(figsize=(8, 6))
for i, loan_amount in enumerate(loan_amounts):
    plt.plot(loan_terms, monthly_payments[0][i], label=f"Loan Amount: ${loan_amount:,}")

plt.title('Monthly Payments for Different Loan Amounts and Terms')
plt.xlabel('Loan Term (Years)')
plt.ylabel('Monthly Payment ($)')
plt.legend(title="Loan Amount")
plt.tight_layout()

# Save the chart as an image (relative path)
chart_path = 'monthly_payments_chart.png'
plt.savefig(chart_path)
plt.close()

# Add the chart to the slide
slide_3.shapes.add_picture(chart_path, Inches(1), Inches(1), width=Inches(8))

# Save the PowerPoint presentation
file_path = 'Financial_Data_Analysis_Presentation.pptx'
prs.save(file_path)

print(f"PowerPoint presentation saved as {file_path}")
